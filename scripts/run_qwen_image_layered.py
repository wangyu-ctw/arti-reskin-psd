#!/usr/bin/env python3
"""Run Qwen-Image-Layered from a local model directory.

Example:
  python scripts/run_qwen_image_layered.py \
    --prompt "Describe the full image content here" \
    --input-path /workspace/input/binan/text.png \
    --output-path /workspace/output/QwenImageLayered \
    --model-dir /workspace/hf_models/Qwen-Image-Layered
"""

from __future__ import annotations

import argparse
import json
import os
import random
import sys
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Iterable

os.environ.setdefault("PYTORCH_CUDA_ALLOC_CONF", "expandable_segments:True")


MAX_SEED = 2**31 - 1
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
QWEN_LAYERED_REQUIRED_PATHS = (
    "model_index.json",
    "processor",
    "scheduler",
    "text_encoder/config.json",
    "text_encoder/model.safetensors.index.json",
    "text_encoder/model-00001-of-00004.safetensors",
    "tokenizer",
    "transformer/config.json",
    "transformer/diffusion_pytorch_model.safetensors.index.json",
    "transformer/diffusion_pytorch_model-00001-of-00005.safetensors",
    "vae/config.json",
    "vae/diffusion_pytorch_model.safetensors",
)
QWEN_LAYERED_LARGE_FILES = (
    "text_encoder/model-00001-of-00004.safetensors",
    "transformer/diffusion_pytorch_model-00001-of-00005.safetensors",
    "vae/diffusion_pytorch_model.safetensors",
)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Run Qwen-Image-Layered and save generated layers.",
        formatter_class=argparse.ArgumentDefaultsHelpFormatter,
    )
    parser.add_argument(
        "--model-dir",
        default="/workspace/hf_models/Qwen-Image-Layered",
        help="Local Qwen-Image-Layered model directory.",
    )
    parser.add_argument(
        "--input-path",
        "--input-image",
        dest="input_path",
        default="/workspace/input/binan/text.png",
        help="Input image file or a directory containing images.",
    )
    parser.add_argument(
        "--recursive-input",
        action="store_true",
        help="When --input-path is a directory, scan images recursively.",
    )
    parser.add_argument(
        "--output-path",
        "--output-root",
        "--output-dir",
        dest="output_path",
        default="/workspace/output/QwenImageLayered",
        help="Output base directory. A timestamped run directory is created under this path.",
    )
    parser.add_argument(
        "--run-name",
        default=None,
        help="Run directory name under --output-path. Defaults to a timestamp.",
    )
    parser.add_argument("--prompt", default="", help="Prompt describing the input image.")
    parser.add_argument(
        "--prompt-file",
        default=None,
        help="Read prompt text from this file. Overrides --prompt when provided.",
    )
    parser.add_argument("--negative-prompt", default=" ", help="Negative prompt.")
    parser.add_argument("--seed", type=int, default=777, help="Random seed.")
    parser.add_argument(
        "--randomize-seed",
        action="store_true",
        help="Ignore --seed and generate a random seed.",
    )
    parser.add_argument(
        "--true-cfg-scale",
        type=float,
        default=4.0,
        help="True CFG scale.",
    )
    parser.add_argument(
        "--steps",
        type=int,
        default=50,
        help="Number of inference steps.",
    )
    parser.add_argument(
        "--layers",
        type=int,
        default=4,
        help="Number of output layers.",
    )
    parser.add_argument(
        "--resolution",
        type=int,
        default=640,
        help="Generation bucket resolution. The official app recommends 640.",
    )
    parser.add_argument(
        "--cfg-normalize",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Enable CFG normalization.",
    )
    parser.add_argument(
        "--use-en-prompt",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Use English automatic captioning when no prompt is provided.",
    )
    parser.add_argument(
        "--device",
        default="cuda",
        help="Torch device for generation.",
    )
    parser.add_argument(
        "--gpu-id",
        type=int,
        default=0,
        help="GPU id used by CPU offload modes.",
    )
    parser.add_argument(
        "--dtype",
        choices=("bfloat16", "float16", "float32"),
        default="bfloat16",
        help="Torch dtype for loading/running the model.",
    )
    parser.add_argument(
        "--load-mode",
        choices=("cuda", "model-offload", "sequential-offload"),
        default="cuda",
        help=(
            "Model placement mode. Use sequential-offload for 24GB GPUs; "
            "cuda is fastest but needs much more VRAM."
        ),
    )
    parser.add_argument(
        "--local-files-only",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Load only from local model files.",
    )
    parser.add_argument(
        "--low-cpu-mem-usage",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Use lower CPU memory loading path when supported.",
    )
    parser.add_argument(
        "--model-file-check",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Validate the expected Qwen-Image-Layered local file layout before loading.",
    )
    parser.add_argument(
        "--vae-tiling",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Enable VAE tiling if supported.",
    )
    parser.add_argument(
        "--vae-slicing",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Enable VAE slicing if supported.",
    )
    parser.add_argument(
        "--attention-slicing",
        action=argparse.BooleanOptionalAction,
        default=False,
        help="Enable attention slicing if supported.",
    )
    parser.add_argument(
        "--save-composite",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Save a simple alpha-composited preview.",
    )
    parser.add_argument("--save-zip", action="store_true", help="Save layers.zip.")
    parser.add_argument("--save-pptx", action="store_true", help="Save layers.pptx.")
    parser.add_argument("--save-psd", action="store_true", help="Save layers.psd.")
    return parser.parse_args()


def torch_dtype(torch_module: object, name: str) -> object:
    if name == "bfloat16":
        return torch_module.bfloat16
    if name == "float16":
        return torch_module.float16
    if name == "float32":
        return torch_module.float32
    raise ValueError(f"Unsupported dtype: {name}")


def read_prompt(args: argparse.Namespace) -> str:
    if args.prompt_file:
        return Path(args.prompt_file).read_text(encoding="utf-8").strip()
    return args.prompt


def make_run_output_dir(args: argparse.Namespace) -> Path:
    run_name = args.run_name or datetime.now().strftime("%Y%m%d_%H%M%S")
    output_dir = Path(args.output_path) / run_name
    output_dir.mkdir(parents=True, exist_ok=True)
    return output_dir


def resolve_input_images(input_path: Path, recursive: bool) -> list[Path]:
    if input_path.is_file():
        if input_path.suffix.lower() not in IMAGE_SUFFIXES:
            raise ValueError(f"Input file is not a supported image: {input_path}")
        return [input_path]

    if not input_path.is_dir():
        raise FileNotFoundError(f"Input path does not exist: {input_path}")

    pattern_iter = input_path.rglob("*") if recursive else input_path.iterdir()
    images = sorted(
        path
        for path in pattern_iter
        if path.is_file() and path.suffix.lower() in IMAGE_SUFFIXES
    )
    if not images:
        raise ValueError(f"No supported images found in input directory: {input_path}")
    return images


def safe_dir_name(path: Path, used_names: set[str]) -> str:
    cleaned = "".join(ch if ch.isalnum() or ch in "._-" else "_" for ch in path.stem)
    cleaned = cleaned.strip("._-") or "image"
    candidate = cleaned
    counter = 2
    while candidate in used_names:
        candidate = f"{cleaned}_{counter:02d}"
        counter += 1
    used_names.add(candidate)
    return candidate


def validate_qwen_layered_model_dir(model_dir: Path) -> list[str]:
    problems = []
    for relative_path in QWEN_LAYERED_REQUIRED_PATHS:
        path = model_dir / relative_path
        if not path.exists():
            problems.append(f"missing: {relative_path}")

    for relative_path in QWEN_LAYERED_LARGE_FILES:
        path = model_dir / relative_path
        if path.exists() and path.is_file() and path.stat().st_size < 1024 * 1024:
            problems.append(
                f"too small, likely a Git LFS/Xet pointer instead of real weights: {relative_path}"
            )
    return problems


def print_model_download_hint(model_dir: Path) -> None:
    print("", file=sys.stderr)
    print("The local model directory does not look like a complete Qwen/Qwen-Image-Layered download.", file=sys.stderr)
    print("Re-download or resume it with:", file=sys.stderr)
    print("", file=sys.stderr)
    print('  pip install -U "huggingface_hub[cli,hf_xet]"', file=sys.stderr)
    print(f"  hf download Qwen/Qwen-Image-Layered --local-dir {model_dir}", file=sys.stderr)
    print("", file=sys.stderr)
    print("If your environment only has huggingface-cli, use:", file=sys.stderr)
    print(f"  huggingface-cli download Qwen/Qwen-Image-Layered --local-dir {model_dir}", file=sys.stderr)
    print("", file=sys.stderr)


def load_input_image(path: str) -> object:
    from PIL import Image

    return Image.open(path).convert("RGB").convert("RGBA")


def normalize_output_images(raw_images: object) -> list[object]:
    from PIL import Image

    if not isinstance(raw_images, list):
        raise TypeError(f"Expected output.images to be a list, got {type(raw_images)!r}")
    if not raw_images:
        raise ValueError("Pipeline returned no images.")

    first = raw_images[0]
    if isinstance(first, Image.Image):
        return raw_images
    if isinstance(first, list) and first and isinstance(first[0], Image.Image):
        return first
    raise TypeError(f"Unexpected output image structure: {type(first)!r}")


def enable_if_supported(pipe: object, method_name: str) -> None:
    method = getattr(pipe, method_name, None)
    if callable(method):
        method()


def place_pipeline(pipe: object, args: argparse.Namespace, torch_module: object) -> object:
    if args.load_mode == "cuda":
        return pipe.to(args.device, torch_dtype(torch_module, args.dtype))
    if args.load_mode == "model-offload":
        pipe.enable_model_cpu_offload(gpu_id=args.gpu_id)
        return pipe
    if args.load_mode == "sequential-offload":
        pipe.enable_sequential_cpu_offload(gpu_id=args.gpu_id)
        return pipe
    raise ValueError(f"Unsupported load mode: {args.load_mode}")


def save_layers(layers: Iterable[object], output_dir: Path) -> list[Path]:
    layer_paths = []
    for idx, layer in enumerate(layers, start=1):
        path = output_dir / f"layer_{idx:02d}.png"
        layer.save(path)
        layer_paths.append(path)
    return layer_paths


def save_composite(layers: list[object], output_dir: Path) -> Path:
    from PIL import Image

    canvas = Image.new("RGBA", layers[0].size, (0, 0, 0, 0))
    for layer in layers:
        canvas.alpha_composite(layer.convert("RGBA"))
    path = output_dir / "composite.png"
    canvas.save(path)
    return path


def save_zip(layer_paths: list[Path], output_dir: Path) -> Path:
    zip_path = output_dir / "layers.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zipf:
        for path in layer_paths:
            zipf.write(path, arcname=path.name)
    return zip_path


def save_pptx(layer_paths: list[Path], output_dir: Path) -> Path:
    from PIL import Image
    from pptx import Presentation

    with Image.open(layer_paths[0]) as img:
        img_width_px, img_height_px = img.size

    def px_to_emu(px: int, dpi: int = 96) -> int:
        return int((px / dpi) * 914400)

    prs = Presentation()
    prs.slide_width = px_to_emu(img_width_px)
    prs.slide_height = px_to_emu(img_height_px)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for path in layer_paths:
        slide.shapes.add_picture(
            str(path),
            0,
            0,
            width=px_to_emu(img_width_px),
            height=px_to_emu(img_height_px),
        )

    pptx_path = output_dir / "layers.pptx"
    prs.save(pptx_path)
    return pptx_path


def save_psd(layer_paths: list[Path], output_dir: Path) -> Path:
    from PIL import Image
    from psd_tools import PSDImage

    layers = [Image.open(path).convert("RGBA") for path in layer_paths]
    width, height = layers[0].size
    psd = PSDImage.new(mode="RGBA", size=(width, height))
    for idx, image in enumerate(layers, start=1):
        psd.append(psd.create_pixel_layer(image=image, name=f"Layer {idx}"))

    psd_path = output_dir / "layers.psd"
    psd.save(psd_path)
    for image in layers:
        image.close()
    return psd_path


def main() -> int:
    args = parse_args()

    model_dir = Path(args.model_dir)
    input_path = Path(args.input_path)
    output_path = Path(args.output_path)
    if not model_dir.exists():
        print(f"Model directory does not exist: {model_dir}", file=sys.stderr)
        return 2
    if args.model_file_check:
        model_problems = validate_qwen_layered_model_dir(model_dir)
        if model_problems:
            print(f"Model directory is incomplete or invalid: {model_dir}", file=sys.stderr)
            for problem in model_problems:
                print(f"  - {problem}", file=sys.stderr)
            print_model_download_hint(model_dir)
            return 2
    if output_path.exists() and not output_path.is_dir():
        print(f"Output path exists but is not a directory: {output_path}", file=sys.stderr)
        return 2
    if args.layers < 1:
        print("--layers must be >= 1", file=sys.stderr)
        return 2
    if args.steps < 1:
        print("--steps must be >= 1", file=sys.stderr)
        return 2
    try:
        input_images = resolve_input_images(input_path, args.recursive_input)
    except (OSError, ValueError) as exc:
        print(str(exc), file=sys.stderr)
        return 2

    try:
        import torch
        from diffusers import DiffusionPipeline
    except ImportError as exc:
        print(f"Missing runtime dependency: {exc}", file=sys.stderr)
        print(
            "Install torch, pillow, diffusers, transformers, accelerate, and optional export packages.",
            file=sys.stderr,
        )
        return 2

    if args.load_mode == "cuda" and args.device.startswith("cuda") and not torch.cuda.is_available():
        print("CUDA is not available. Use --device cpu or an offload mode.", file=sys.stderr)
        return 2

    prompt = read_prompt(args)
    run_output_dir = make_run_output_dir(args)
    batch_mode = len(input_images) > 1 or input_path.is_dir()
    used_output_names: set[str] = set()

    print(f"Loading model: {model_dir}")
    print(f"Load mode: {args.load_mode}, dtype: {args.dtype}")
    print(f"Input path: {input_path}")
    print(f"Images found: {len(input_images)}")
    print(f"Run output dir: {run_output_dir}")
    pipe = DiffusionPipeline.from_pretrained(
        str(model_dir),
        torch_dtype=torch_dtype(torch, args.dtype),
        local_files_only=args.local_files_only,
        low_cpu_mem_usage=args.low_cpu_mem_usage,
    )
    pipe = place_pipeline(pipe, args, torch)
    pipe.set_progress_bar_config(disable=False)

    if args.vae_tiling:
        enable_if_supported(pipe, "enable_vae_tiling")
    if args.vae_slicing:
        enable_if_supported(pipe, "enable_vae_slicing")
    if args.attention_slicing:
        enable_if_supported(pipe, "enable_attention_slicing")

    generator_device = "cuda" if args.device.startswith("cuda") and torch.cuda.is_available() else "cpu"

    shared_metadata = {
        "model_dir": str(model_dir),
        "input_path": str(input_path),
        "output_path": str(output_path),
        "run_output_dir": str(run_output_dir),
        "prompt": prompt,
        "negative_prompt": args.negative_prompt,
        "true_cfg_scale": args.true_cfg_scale,
        "steps": args.steps,
        "layers": args.layers,
        "resolution": args.resolution,
        "cfg_normalize": args.cfg_normalize,
        "use_en_prompt": args.use_en_prompt,
        "dtype": args.dtype,
        "load_mode": args.load_mode,
        "argv": sys.argv,
    }

    manifest = {
        **shared_metadata,
        "image_count": len(input_images),
        "items": [],
    }

    for index, image_path in enumerate(input_images, start=1):
        seed = random.randint(0, MAX_SEED) if args.randomize_seed else args.seed
        item_output_dir = run_output_dir
        if batch_mode:
            item_output_dir = run_output_dir / safe_dir_name(image_path, used_output_names)
            item_output_dir.mkdir(parents=True, exist_ok=True)

        generator = torch.Generator(device=generator_device).manual_seed(seed)
        inputs = {
            "image": load_input_image(str(image_path)),
            "generator": generator,
            "true_cfg_scale": args.true_cfg_scale,
            "prompt": prompt,
            "negative_prompt": args.negative_prompt,
            "num_inference_steps": args.steps,
            "num_images_per_prompt": 1,
            "layers": args.layers,
            "resolution": args.resolution,
            "cfg_normalize": args.cfg_normalize,
            "use_en_prompt": args.use_en_prompt,
        }

        print(f"[{index}/{len(input_images)}] Input image: {image_path}")
        print(f"[{index}/{len(input_images)}] Output dir: {item_output_dir}")
        print(f"[{index}/{len(input_images)}] Seed: {seed}")
        with torch.inference_mode():
            output = pipe(**inputs)

        output_layers = normalize_output_images(output.images)
        layer_paths = save_layers(output_layers, item_output_dir)
        saved = {"layers": [str(path) for path in layer_paths]}

        if args.save_composite:
            saved["composite"] = str(save_composite(output_layers, item_output_dir))
        if args.save_zip:
            saved["zip"] = str(save_zip(layer_paths, item_output_dir))
        if args.save_pptx:
            saved["pptx"] = str(save_pptx(layer_paths, item_output_dir))
        if args.save_psd:
            saved["psd"] = str(save_psd(layer_paths, item_output_dir))

        item_metadata = {
            **shared_metadata,
            "input_image": str(image_path),
            "output_dir": str(item_output_dir),
            "seed": seed,
            "saved": saved,
        }
        metadata_path = item_output_dir / "metadata.json"
        metadata_path.write_text(
            json.dumps(item_metadata, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

        manifest["items"].append(
            {
                "input_image": str(image_path),
                "output_dir": str(item_output_dir),
                "seed": seed,
                "saved": saved,
                "metadata": str(metadata_path),
            }
        )

        print("Saved files:")
        for key, value in saved.items():
            print(f"  {key}: {value}")
        print(f"  metadata: {metadata_path}")

    manifest_path = run_output_dir / "manifest.json"
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    print(f"Run manifest: {manifest_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
